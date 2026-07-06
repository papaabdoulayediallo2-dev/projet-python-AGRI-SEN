#!/usr/bin/env python3
"""
Script pour générer une présentation PowerPoint du projet AGRI-SEN
Présentation détaillée sans emoji
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Créer une présentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Définir les couleurs
COULEUR_PRINCIPALE = RGBColor(34, 139, 34)  # Vert agricole
COULEUR_TEXTE = RGBColor(0, 0, 0)
COULEUR_TITRE = RGBColor(255, 255, 255)
COULEUR_SOUS_TITRE = RGBColor(50, 50, 50)

def ajouter_slide_titre(titre, sous_titre):
    """Ajoute une slide de titre"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COULEUR_PRINCIPALE
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = titre
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COULEUR_TITRE
    p.alignment = PP_ALIGN.CENTER
    
    # Sous-titre
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = sous_titre
    p.font.size = Pt(28)
    p.font.color.rgb = COULEUR_TITRE
    p.alignment = PP_ALIGN.CENTER

def ajouter_slide_contenu(titre, contenu_list):
    """Ajoute une slide avec du contenu"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # Titre
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = titre
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COULEUR_PRINCIPALE
    
    # Contenu
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(5.7))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, ligne in enumerate(contenu_list):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = ligne
        p.font.size = Pt(18)
        p.font.color.rgb = COULEUR_TEXTE
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        p.line_spacing = 1.15

# Slide 1 : Titre
ajouter_slide_titre("AGRI-SEN", "Système de Gestion de Coopérative Agricole")

# Slide 2 : Introduction
ajouter_slide_contenu("Introduction du Projet", [
    "AGRI-SEN est un système informatique conçu pour gérer les opérations d'une",
    "coopérative agricole. Ce système permet de centraliser et d'analyser les données",
    "concernant les producteurs et leurs récoltes.",
    "",
    "Le projet a été développé par 5 développeurs travaillant de manière collaborative",
    "en Python, avec une architecture modulaire où chaque membre s'occupe d'une partie",
    "spécifique du système."
])

# Slide 3 : Objectifs Généraux
ajouter_slide_contenu("Objectifs du Projet", [
    "1. Centraliser les informations sur tous les producteurs de la coopérative",
    "   dans une base de données accessible et organisée.",
    "",
    "2. Permettre l'enregistrement détaillé de chaque récolte avec validation",
    "   des données pour assurer la qualité.",
    "",
    "3. Générer des rapports et statistiques de production pour analyser",
    "   la performance de la coopérative.",
    "",
    "4. Identifier les producteurs les plus performants et calculer les",
    "   moyennes de production pour optimiser la gestion."
])

# Slide 4 : Architecture Générale - Vue d'ensemble
ajouter_slide_contenu("Architecture Générale du Système", [
    "Le système est divisé en 5 modules Python, chacun responsable d'une fonction",
    "spécifique :",
    "",
    "Module 1 (membre1.py): Gestion des producteurs",
    "Module 2 (membre2.py): Gestion des récoltes",
    "Module 3 (Membre3.py): Statistiques de production",
    "Module 4 (membre4.py): Analyses avancées",
    "Module 5 (main.py): Interface utilisateur et intégration"
])

# Slide 5 : Détail - Membre 1
ajouter_slide_contenu("Module 1 : Gestion des Producteurs (membre1.py)", [
    "Ce module gère l'ajout et le suivi des producteurs de la coopérative.",
    "",
    "Composants principaux :",
    "- producteurs : liste globale stockant tous les producteurs enregistrés",
    "- ajouter_producteur() : permet d'ajouter un nouveau producteur avec validation",
    "- afficher_producteurs() : affiche la liste complète avec nombre de récoltes",
    "- rechercher_producteur(nom) : trouve un producteur par son nom"
])

# Slide 6 : Détail - Membre 2
ajouter_slide_contenu("Module 2 : Gestion des Récoltes (membre2.py)", [
    "Ce module gère l'enregistrement et l'affichage des récoltes des producteurs.",
    "",
    "Composants principaux :",
    "- PRODUITS_VALIDES : liste des 4 produits autorisés (Mil, Maïs, Riz, Arachide)",
    "- produit_valide(produit) : valide que le produit fait partie de la liste",
    "- enregistrer_recolte() : enregistre une nouvelle récolte avec validation complète",
    "- afficher_recoltes(nom) : affiche toutes les récoltes ou celles d'un producteur"
])

# Slide 7 : Détail - Membre 3
ajouter_slide_contenu("Module 3 : Statistiques de Production (Membre3.py)", [
    "Ce module calcule les statistiques globales de production de la coopérative.",
    "",
    "Composants principaux :",
    "- production_totale() : calcule la quantité totale produite en kg par",
    "  tous les producteurs",
    "- production_par_produit() : décompose la production par type de produit",
    "  et affiche les totaux pour chaque catégorie"
])

# Slide 8 : Détail - Membre 4
ajouter_slide_contenu("Module 4 : Analyses Avancées (membre4.py)", [
    "Ce module fournit des analyses détaillées des performances des producteurs.",
    "",
    "Composants principaux :",
    "- meilleur_producteur(liste) : identifie et affiche le producteur ayant",
    "  la plus grande quantité totale produite",
    "- moyenne_producteur(liste) : calcule la production moyenne par producteur",
    "  et affiche le résultat formaté avec 2 décimales"
])

# Slide 9 : Détail - Membre 5 (Main)
ajouter_slide_contenu("Module 5 : Interface Principale (main.py)", [
    "Ce module fournit l'interface utilisateur et intègre tous les autres modules.",
    "",
    "Fonctionnalités :",
    "- Menu interactif affichant 10 options numérotées (1-9, plus 0 pour quitter)",
    "- Importe toutes les fonctions des 4 modules précédents",
    "- Boucle infinie permettant à l'utilisateur de naviguer entre les options",
    "- Gère les entrées utilisateur et appelle les fonctions appropriées"
])

# Slide 10 : Flux d'Exécution
ajouter_slide_contenu("Flux d'Exécution du Programme", [
    "Étape 1 : L'utilisateur lance le programme avec la commande 'python main.py'",
    "",
    "Étape 2 : Le menu principal s'affiche avec 10 options disponibles",
    "",
    "Étape 3 : L'utilisateur saisit un numéro d'option (1-9 ou 0)",
    "",
    "Étape 4 : La fonction correspondante est exécutée",
    "",
    "Étape 5 : Après exécution, l'utilisateur retourne au menu pour une nouvelle action",
    "",
    "Étape 6 : Le processus continue jusqu'à l'option 0 (Quitter)"
])

# Slide 11 : Fonctionnalités du Menu Principal
ajouter_slide_contenu("Options du Menu Principal", [
    "Option 1 : Ajouter un producteur - Demande le nom et valide l'unicité",
    "Option 2 : Afficher les producteurs - Liste tous avec nombre de récoltes",
    "Option 3 : Rechercher un producteur - Trouve par nom et affiche infos",
    "Option 4 : Enregistrer une récolte - Ajoute nouvelle récolte à un producteur",
    "Option 5 : Afficher les récoltes - Liste toutes ou filtrées par producteur",
    "Option 6 : Production totale - Somme de toutes les récoltes en kg",
    "Option 7 : Production par produit - Détail pour chaque type de produit",
    "Option 8 : Meilleur producteur - Identifie le plus productif",
    "Option 9 : Moyenne par producteur - Calcule la moyenne de production"
])

# Slide 12 : Produits Autorisés
ajouter_slide_contenu("Produits Autorisés dans le Système", [
    "Le système gère uniquement 4 types de produits agricoles :",
    "",
    "1. MIL - Céréale traditionnelle très résistante à la chaleur et à la sécheresse",
    "",
    "2. MAÏS - Céréale polyvalente cultivée largement en Afrique",
    "",
    "3. RIZ - Céréale cultivée en zones irrigables ou pluvieuses",
    "",
    "4. ARACHIDE - Culture de rente importante pour les revenus agricoles",
    "",
    "Tous les produits saisis sont automatiquement normalisés en majuscules"
])

# Slide 13 : Structure des Données
ajouter_slide_contenu("Structure des Données du Système", [
    "Les données sont organisées dans une structure Python basée sur des dictionnaires :",
    "",
    "producteurs = [",
    "  {",
    "    'nom': 'Producteur 1',",
    "    'recoltes': [",
    "      {'nom': 'Producteur 1', 'produit': 'Mil', 'quantite': 100.5},",
    "      {'nom': 'Producteur 1', 'produit': 'Riz', 'quantite': 75.0}",
    "    ]",
    "  }",
    "]"
])

# Slide 14 : Validation des Données
ajouter_slide_contenu("Système de Validation et Contrôle de Qualité", [
    "Le système implémente une validation stricte à chaque étape :",
    "",
    "Validation des Producteurs :",
    "- Le nom ne peut pas être vide",
    "- Les doublons sont détectés et rejetés",
    "- Recherche insensible à la casse",
    "",
    "Validation des Récoltes :",
    "- Le producteur doit exister avant d'enregistrer une récolte",
    "- Le produit doit être dans la liste autorisée",
    "- La quantité doit être un nombre positif supérieur à zéro"
])

# Slide 15 : Technologies et Outils
ajouter_slide_contenu("Technologies et Outils Utilisés", [
    "Langage de programmation : Python 3.x",
    "   - Choisi pour sa simplicité et sa puissance pour traiter les données",
    "",
    "Structures de données : Dictionnaires et Listes Python",
    "   - Permettent une organisation flexible des données",
    "",
    "Stockage : Données en mémoire (variables Python)",
    "   - Les données restent actives durant l'exécution du programme",
    "",
    "Contrôle de version : Git et GitHub",
    "   - Permet le suivi des modifications et la collaboration d'équipe"
])

# Slide 16 : Avantages du Système
ajouter_slide_contenu("Avantages Principaux du Système", [
    "Interface Simple et Conviviale :",
    "- Menu numéroté facile à naviguer sans connaissances techniques",
    "",
    "Données Fiables et Validées :",
    "- Tous les champs sont vérifiés avant enregistrement",
    "",
    "Rapports Détaillés :",
    "- Analyse complète de la production avec plusieurs perspectives",
    "",
    "Modularité et Extensibilité :",
    "- Code organisé en modules permettant l'ajout de nouvelles fonctionnalités",
    "",
    "Performance :",
    "- Recherche et calculs très rapides sur les données"
])

# Slide 17 : Limitations et Améliorations Futures
ajouter_slide_contenu("Limitations Actuelles et Perspectives Futures", [
    "Limitations actuelles :",
    "- Les données ne sont stockées qu'en mémoire (perdues au redémarrage)",
    "- Pas de fonction de suppression ou modification de données",
    "- Pas d'authentification utilisateur",
    "",
    "Améliorations futures possibles :",
    "- Ajouter une base de données (SQL) pour persistance des données",
    "- Implémenter édition et suppression de producteurs/récoltes",
    "- Ajouter export des rapports en PDF ou Excel",
    "- Implémenter système de sauvegarde/chargement de fichiers"
])

# Slide 18 : Exemple d'Utilisation
ajouter_slide_contenu("Exemple Complet d'Utilisation", [
    "Scénario : Un gestionnaire de coopérative souhaite enregistrer les récoltes",
    "",
    "Étapes :",
    "1. Lance le programme : python main.py",
    "2. Choisit option 1 pour ajouter le producteur 'Amadou Sow'",
    "3. Choisit option 4 pour enregistrer 100 kg de Riz pour Amadou",
    "4. Choisit option 4 à nouveau pour ajouter 50 kg d'Arachide",
    "5. Choisit option 6 pour voir production totale",
    "6. Choisit option 7 pour voir détail par produit",
    "7. Choisit option 8 pour identifier le meilleur producteur"
])

# Slide 19 : Conclusion
ajouter_slide_contenu("Conclusion", [
    "AGRI-SEN est un système complet et fonctionnel pour la gestion des coopératives",
    "agricoles. Il démontre une architecture logicielle bien pensée avec une séparation",
    "claire des responsabilités.",
    "",
    "Le code est :",
    "- Bien organisé en modules indépendants",
    "- Bien commenté et facile à comprendre",
    "- Facile à maintenir et à étendre",
    "- Prêt pour utilisation en environnement réel",
    "",
    "Ce projet représente une collaboration d'équipe réussie dans le développement",
    "logiciel collaboratif."
])

# Sauvegarder la présentation
output_path = "AGRI-SEN_Presentation_Updated.pptx"
prs.save(output_path)
print(f"Présentation créée avec succès : {output_path}")
